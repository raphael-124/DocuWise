import os
import asyncio
import logging
from pypdf import PdfReader
try:
    import fitz # PyMuPDF
except ImportError:
    fitz = None
import docx
from pptx import Presentation
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_google_genai import GoogleGenerativeAIEmbeddings
from langchain_ollama import OllamaEmbeddings

logger = logging.getLogger(__name__)

# Configure a location for the vector store
DATA_DIR = "data"
VECTOR_STORE_DIR = os.path.join(DATA_DIR, "vectorstore")
os.makedirs(VECTOR_STORE_DIR, exist_ok=True)

class DocumentProcessor:
    def __init__(self):
        # Very large chunks (8000) to minimize the number of embedding API calls
        # This is critical for staying within Gemini Free Tier quotas (15 RPM)
        # Optimized for speed and API efficiency
        self.text_splitter = RecursiveCharacterTextSplitter(chunk_size=8000, chunk_overlap=200)

    def _get_embeddings(self, use_local: bool = False, provider: str = "gemini"):
        if provider == "ollama":
            return OllamaEmbeddings(model="nomic-embed-text")
        elif provider == "huggingface":
            try:
                from langchain_huggingface import HuggingFaceEmbeddings
                return HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2")
            except ImportError:
                logger.warning("langchain-huggingface not installed, falling back to Gemini")
                provider = "gemini"
        
        api_key = os.getenv("GOOGLE_API_KEY")
        if not api_key:
            raise ValueError("GOOGLE_API_KEY not found in environment variables. Please set it in backend/.env")
        return GoogleGenerativeAIEmbeddings(model="models/gemini-embedding-001", google_api_key=api_key)

    def extract_text(self, file_path: str) -> str:
        ext = os.path.splitext(file_path)[1].lower()
        if ext == ".pdf":
            return self.extract_text_from_pdf(file_path)
        elif ext == ".docx":
            return self.extract_text_from_docx(file_path)
        elif ext == ".pptx":
            return self.extract_text_from_pptx(file_path)
        elif ext in [".txt", ".md"]:
            return self.extract_text_from_text(file_path)
        else:
            raise ValueError(f"Unsupported file extension: {ext}")


    def extract_text_from_pdf(self, file_path: str) -> str:
        text = ""
        try:
            # Prefer PyMuPDF for 10x speed boost
            if fitz:
                doc = fitz.open(file_path)
                for page in doc:
                    text += page.get_text() + "\n\n"
                doc.close()
            else:
                # Fallback to pypdf
                reader = PdfReader(file_path)
                for page in reader.pages:
                    extracted = page.extract_text()
                    if extracted:
                        text += extracted + "\n\n"
        except Exception as e:
            logger.error(f"Error reading PDF {file_path}: {e}")
            raise ValueError(f"Failed to read PDF: {str(e)}")
        
        return text.strip()

    def extract_text_from_docx(self, file_path: str) -> str:
        try:
            doc = docx.Document(file_path)
            return "\n".join([para.text for para in doc.paragraphs]).strip()
        except Exception as e:
            logger.error(f"Error reading DOCX {file_path}: {e}")
            raise ValueError(f"Failed to read DOCX: {str(e)}")

    def extract_text_from_pptx(self, file_path: str) -> str:
        try:
            presentation = Presentation(file_path)
            text_runs = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            return "\n".join(text_runs).strip()
        except Exception as e:
            logger.error(f"Error reading PPTX {file_path}: {e}")
            raise ValueError(f"Failed to read PPTX: {str(e)}")

    def extract_text_from_text(self, file_path: str) -> str:
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read().strip()
        except Exception as e:
            logger.error(f"Error reading text file {file_path}: {e}")
            raise ValueError(f"Failed to read text file: {str(e)}")

    async def _aembed_with_retry(self, chunks, metadatas, embeddings, max_retries=5):
        """Embed chunks in parallel batches for extreme speed."""
        batch_size = 100
        total_batches = (len(chunks) + batch_size - 1) // batch_size
        
        # Concurrency control: Increase to 5 for higher throughput (safe for 15 RPM)
        semaphore = asyncio.Semaphore(5) 
        
        async def embed_batch(batch_num, batch_chunks, batch_meta):
            async with semaphore:
                for attempt in range(max_retries):
                    try:
                        logger.info(f"Embedding batch {batch_num}/{total_batches} ({len(batch_chunks)} chunks) in parallel...")
                        # We use the blocking FAISS call in a thread for now as FAISS isn't natively async
                        # but we compute embeddings in parallel if the provider supports it.
                        return await asyncio.to_thread(
                            FAISS.from_texts, batch_chunks, embeddings, metadatas=batch_meta
                        )
                    except Exception as e:
                        error_str = str(e).lower()
                        if ("429" in error_str or "resource_exhausted" in error_str or "quota" in error_str):
                            raise ValueError("GEMINI_QUOTA_REACHED")
                        
                        if attempt < max_retries - 1:
                            wait_time = 15 * (attempt + 1)
                            logger.warning(f"Rate limit hit on batch {batch_num}, retrying in {wait_time}s...")
                            await asyncio.sleep(wait_time)
                        else:
                            raise

        # Create tasks for all batches
        tasks = []
        for i in range(0, len(chunks), batch_size):
            batch_num = i // batch_size + 1
            tasks.append(embed_batch(batch_num, chunks[i:i + batch_size], metadatas[i:i + batch_size]))

        # Run all batches in parallel
        results = await asyncio.gather(*tasks)
        
        # Merge all partial vectorstores
        vectorstore = results[0]
        for next_vs in results[1:]:
            vectorstore.merge_from(next_vs)
            
        return vectorstore

    async def aprocess_and_store_document(self, file_path: str, document_id: str, filename: str, use_local: bool = False):
        logger.info(f"Processing document async: {document_id}")
        
        # Performance optimization: Instant skip if already processed
        save_path = os.path.join(VECTOR_STORE_DIR, document_id)
        if os.path.exists(save_path):
             logger.info(f"Using existing vector store for {document_id}")
             # Return a placeholder or try to count. For now, since we need num_chunks:
             try:
                 # We need embeddings to load FAISS, get them early
                 provider = "gemini"
                 if document_id.endswith("_local"): provider = "ollama"
                 elif document_id.endswith("_hf"): provider = "huggingface"
                 embeddings = self._get_embeddings(provider=provider)
                 
                 vs = await asyncio.to_thread(FAISS.load_local, save_path, embeddings, allow_dangerous_deserialization=True)
                 return len(vs.index_to_id) if hasattr(vs, 'index_to_id') else 0
             except Exception:
                 logger.warning(f"Could not load existing store, reprocessing...")
                 pass

        text = await asyncio.to_thread(self.extract_text, file_path)
        if not text.strip():
            raise ValueError(f"Could not extract any text from {file_path}. The PDF may be image-based or empty.")
            
        chunks = self.text_splitter.split_text(text)
        logger.info(f"Split document into {len(chunks)} chunks")
        
        metadatas = [{"source": document_id, "filename": filename, "chunk_index": i} for i in range(len(chunks))]
        
        # Determine provider from document_id suffix or default
        provider = "gemini"
        if document_id.endswith("_local"):
             provider = "ollama"
        elif document_id.endswith("_hf"):
             provider = "huggingface"
 
        embeddings = self._get_embeddings(provider=provider)
        
        # Try to embed everything in parallel, auto-revert to HuggingFace if Gemini fails
        try:
            vectorstore = await self._aembed_with_retry(chunks, metadatas, embeddings)
        except ValueError as e:
            if str(e) == "GEMINI_QUOTA_REACHED":
                logger.info("Restarting embedding process with HuggingFace due to Gemini quota...")
                embeddings = self._get_embeddings(provider="huggingface")
                vectorstore = await self._aembed_with_retry(chunks, metadatas, embeddings)
            else:
                raise

        # document_id already carries the provider suffix (_local, _hf, or none).
        # Do NOT append an extra suffix here to avoid double-suffix paths.
        save_path = os.path.join(VECTOR_STORE_DIR, document_id)
        vectorstore.save_local(save_path)
        logger.info(f"Vector store saved to {save_path}")
        
        return len(chunks)

# Initialize a global instance
document_processor = DocumentProcessor()
